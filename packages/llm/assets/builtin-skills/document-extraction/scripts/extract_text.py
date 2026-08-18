#!/usr/bin/env python3
"""Extract plain UTF-8 text from a document file.

Usage: extract_text.py FILE [--max-chars N]

Writes extracted text to stdout. Exits 0 on success (possibly with empty
output for a genuinely empty document), non-zero with a reason on stderr
when the format cannot be handled at all.

Degrades gracefully: prefers the best available library, falls back to
CLI tools, and finally to dependency-free stdlib parsing so the script
works even on hosts where none of the optional deps are installed:

- pdf:   pypdf -> pdftotext -> (scanned pages) pdftoppm+tesseract OCR
- docx:  python-docx -> stdlib zipfile + XML strip
- pptx:  python-pptx -> stdlib zipfile + XML strip
- xlsx:  openpyxl    -> stdlib zipfile sharedStrings
- text (.md .txt .csv .tsv .json .html .xml .yaml .yml .log):
         utf-8 -> utf-8-sig -> latin-1 decode; HTML/XML tags stripped
- other: printable-run fallback (like `strings`) as a last resort
"""

import re
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path

DEFAULT_MAX_CHARS = 400_000
OCR_MAX_PAGES = 10

TEXT_EXTENSIONS = {
    ".md", ".txt", ".csv", ".tsv", ".json", ".html", ".htm", ".xml",
    ".yaml", ".yml", ".log",
}


def fail(reason: str) -> "NoReturn":  # noqa: F821 - py3.11 has NoReturn in typing only
    print(reason, file=sys.stderr)
    sys.exit(1)


def decode_bytes(raw: bytes) -> str:
    # utf-8-sig first: it decodes plain UTF-8 too and strips a BOM.
    for encoding in ("utf-8-sig", "utf-8"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    # latin-1 never fails; mojibake beats nothing for triage purposes.
    return raw.decode("latin-1")


def strip_markup(xml: str) -> str:
    # Paragraph-ish closers become newlines so prose keeps its shape.
    xml = re.sub(r"</(w:p|a:p|p|div|h[1-6]|li|tr|br)>", "\n", xml)
    xml = re.sub(r"<br\s*/?>", "\n", xml)
    xml = re.sub(r"<script[\s\S]*?</script>", " ", xml, flags=re.I)
    xml = re.sub(r"<style[\s\S]*?</style>", " ", xml, flags=re.I)
    xml = re.sub(r"<[^>]+>", " ", xml)
    xml = re.sub(r"[ \t]+", " ", xml)
    return re.sub(r"\n{3,}", "\n\n", xml).strip()


def extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader  # type: ignore

        reader = PdfReader(str(path))
        text = "\n".join((page.extract_text() or "") for page in reader.pages)
        if text.strip():
            return text
    except Exception:
        pass
    pdftotext = shutil.which("pdftotext")
    if pdftotext:
        result = subprocess.run(
            [pdftotext, "-layout", str(path), "-"],
            capture_output=True,
            timeout=120,
        )
        if result.returncode == 0 and result.stdout.strip():
            return decode_bytes(result.stdout)
    ocr = extract_pdf_ocr(path)
    if ocr is not None:
        return ocr
    fail(
        "PDF text extraction failed: install pypdf (pip) or poppler-utils "
        "(pdftotext), plus tesseract for scanned documents"
    )


def extract_pdf_ocr(path: Path) -> "str | None":
    """OCR image-only PDFs when poppler + tesseract are present."""
    pdftoppm = shutil.which("pdftoppm")
    tesseract = shutil.which("tesseract")
    if not pdftoppm or not tesseract:
        return None
    import tempfile

    with tempfile.TemporaryDirectory() as tmp:
        render = subprocess.run(
            [
                pdftoppm, "-png", "-r", "200",
                "-f", "1", "-l", str(OCR_MAX_PAGES),
                str(path), f"{tmp}/page",
            ],
            capture_output=True,
            timeout=300,
        )
        if render.returncode != 0:
            return None
        pages = sorted(Path(tmp).glob("page*.png"))
        if not pages:
            return None
        chunks = []
        for page in pages:
            result = subprocess.run(
                [tesseract, str(page), "-"],
                capture_output=True,
                timeout=300,
            )
            if result.returncode == 0:
                chunks.append(decode_bytes(result.stdout))
        text = "\n".join(chunks).strip()
        return text or None


def extract_docx(path: Path) -> str:
    try:
        import docx  # type: ignore

        document = docx.Document(str(path))
        parts = [p.text for p in document.paragraphs]
        for table in document.tables:
            for row in table.rows:
                parts.append("\t".join(cell.text for cell in row.cells))
        return "\n".join(parts)
    except Exception:
        # Fall through: the stdlib parser handles files python-docx
        # rejects (e.g. minimal packages without full rels).
        pass
    with zipfile.ZipFile(path) as zf:
        return strip_markup(decode_bytes(zf.read("word/document.xml")))


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore

        parts = []
        for slide in Presentation(str(path)).slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
        return "\n".join(parts)
    except Exception:
        pass
    with zipfile.ZipFile(path) as zf:
        slides = sorted(
            name
            for name in zf.namelist()
            if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
        )
        return "\n\n".join(strip_markup(decode_bytes(zf.read(s))) for s in slides)


def extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore

        workbook = load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet in workbook.worksheets:
            parts.append(f"# {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = ["" if value is None else str(value) for value in row]
                if any(cells):
                    parts.append("\t".join(cells))
        return "\n".join(parts)
    except Exception:
        pass
    with zipfile.ZipFile(path) as zf:
        names = zf.namelist()
        shared: list[str] = []
        if "xl/sharedStrings.xml" in names:
            xml = decode_bytes(zf.read("xl/sharedStrings.xml"))
            shared = [
                m.group(1) for m in re.finditer(r"<t[^>]*>([^<]*)</t>", xml)
            ]
        parts = []
        for sheet in sorted(n for n in names if re.fullmatch(r"xl/worksheets/sheet\d+\.xml", n)):
            xml = decode_bytes(zf.read(sheet))
            for row in re.finditer(r"<row[^>]*>([\s\S]*?)</row>", xml):
                cells = []
                for cell in re.finditer(r"<c([^>]*)>([\s\S]*?)</c>", row.group(1)):
                    attrs, inner = cell.group(1), cell.group(2)
                    value = re.search(r"<(?:v|t)[^>]*>([^<]*)</(?:v|t)>", inner)
                    if not value:
                        continue
                    text = value.group(1)
                    if 't="s"' in attrs and text.isdigit() and int(text) < len(shared):
                        text = shared[int(text)]
                    cells.append(text)
                if cells:
                    parts.append("\t".join(cells))
        return "\n".join(parts)


def extract_text_file(path: Path) -> str:
    text = decode_bytes(path.read_bytes())
    if path.suffix.lower() in {".html", ".htm", ".xml"}:
        return strip_markup(text)
    return text


def extract_strings_fallback(path: Path) -> str:
    raw = path.read_bytes()
    runs = re.findall(rb"[\x20-\x7e]{4,}", raw)
    return "\n".join(run.decode("ascii") for run in runs)


def main() -> None:
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    max_chars = DEFAULT_MAX_CHARS
    for arg in sys.argv[1:]:
        if arg.startswith("--max-chars="):
            max_chars = int(arg.split("=", 1)[1])
    if len(args) != 1:
        fail("usage: extract_text.py FILE [--max-chars=N]")
    path = Path(args[0])
    if not path.is_file():
        fail(f"not a file: {path}")

    ext = path.suffix.lower()
    if ext == ".pdf":
        text = extract_pdf(path)
    elif ext == ".docx":
        text = extract_docx(path)
    elif ext == ".pptx":
        text = extract_pptx(path)
    elif ext == ".xlsx":
        text = extract_xlsx(path)
    elif ext in TEXT_EXTENSIONS:
        text = extract_text_file(path)
    else:
        text = extract_strings_fallback(path)

    sys.stdout.write(text[:max_chars])


if __name__ == "__main__":
    main()
